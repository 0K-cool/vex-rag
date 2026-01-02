"""
Document Loader - Load and parse documents for indexing

Supports:
- Text: .md, .py, .ts, .js, .json, .txt, .sh, .yml, .yaml
- Rich Documents: .pdf, .docx, .pptx
"""

import os
from pathlib import Path
from typing import Dict, Optional, List
from dataclasses import dataclass
import logging

logger = logging.getLogger(__name__)

# Rich document parsers
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


@dataclass
class Document:
    """Represents a loaded document"""
    content: str
    file_path: str
    project: str
    metadata: Dict


class DocumentLoader:
    """Load documents from filesystem"""

    SUPPORTED_EXTENSIONS = {
        # Text files
        '.md', '.py', '.ts', '.js', '.json', '.txt', '.sh', '.yml', '.yaml',
        # Rich documents
        '.pdf', '.docx', '.pptx'
    }

    def __init__(self):
        self.loaded_count = 0

    def load_file(self, file_path: str, project: str) -> Optional[Document]:
        """
        Load a single file

        Args:
            file_path: Absolute path to file
            project: Project name (PAI, VERSANT-IR, VERSANT-ATHENA, etc.)

        Returns:
            Document object or None if file cannot be loaded
        """
        path = Path(file_path)

        # Check file exists
        if not path.exists():
            logger.warning(f"File not found: {file_path}")
            return None

        # Check extension
        if path.suffix not in self.SUPPORTED_EXTENSIONS:
            logger.warning(f"Unsupported file type: {path.suffix}")
            return None

        try:
            # Route to appropriate parser
            if path.suffix == '.pdf':
                content, metadata = self._parse_pdf(file_path)
            elif path.suffix == '.docx':
                content, metadata = self._parse_docx(file_path)
            elif path.suffix == '.pptx':
                content, metadata = self._parse_pptx(file_path)
            else:
                # Text files
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                metadata = {
                    'file_name': path.name,
                    'file_ext': path.suffix,
                    'file_size': len(content),
                    'absolute_path': str(path.absolute()),
                }

            # Skip empty files
            if not content or not content.strip():
                logger.warning(f"Empty file: {file_path}")
                return None

            # Create document
            doc = Document(
                content=content,
                file_path=file_path,
                project=project,
                metadata=metadata
            )

            self.loaded_count += 1
            return doc

        except Exception as e:
            logger.error(f"Error loading {file_path}: {e}")
            return None

    def load_directory(
        self,
        directory: str,
        project: str,
        recursive: bool = True,
        exclude_patterns: Optional[List[str]] = None
    ) -> List[Document]:
        """
        Load all supported files from a directory

        Args:
            directory: Path to directory
            project: Project name
            recursive: Whether to search subdirectories
            exclude_patterns: List of patterns to exclude (e.g., ['node_modules', '.git'])

        Returns:
            List of Document objects
        """
        if exclude_patterns is None:
            exclude_patterns = [
                'node_modules',
                '.git',
                '.venv',
                '.venv-rag',
                '__pycache__',
                'dist',
                'build',
                'output/temp',
                '.DS_Store'
            ]

        documents = []
        dir_path = Path(directory)

        if not dir_path.exists():
            logger.error(f"Directory not found: {directory}")
            return documents

        # Get files
        if recursive:
            files = dir_path.rglob('*')
        else:
            files = dir_path.glob('*')

        for file_path in files:
            # Skip directories
            if file_path.is_dir():
                continue

            # Skip excluded patterns
            if any(pattern in str(file_path) for pattern in exclude_patterns):
                continue

            # Load file
            doc = self.load_file(str(file_path), project)
            if doc:
                documents.append(doc)

        logger.info(f"Loaded {len(documents)} documents from {directory}")
        return documents

    def _parse_pdf(self, file_path: str) -> tuple:
        """
        Parse PDF file using PyMuPDF

        Args:
            file_path: Path to PDF file

        Returns:
            (content, metadata) tuple
        """
        if not PYMUPDF_AVAILABLE:
            raise ImportError("PyMuPDF not installed. Run: pip install pymupdf")

        path = Path(file_path)
        doc = fitz.open(file_path)

        # Extract text from all pages
        text_parts = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{text}")

        content = "\n\n".join(text_parts)

        # Create metadata
        metadata = {
            'file_name': path.name,
            'file_ext': path.suffix,
            'file_size': len(content),
            'absolute_path': str(path.absolute()),
            'page_count': len(doc),
            'document_type': 'pdf'
        }

        doc.close()
        return content, metadata

    def _parse_docx(self, file_path: str) -> tuple:
        """
        Parse Word document using python-docx

        Args:
            file_path: Path to .docx file

        Returns:
            (content, metadata) tuple
        """
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

        path = Path(file_path)
        doc = DocxDocument(file_path)

        # Extract text from all paragraphs
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)

        content = "\n\n".join(paragraphs)

        # Create metadata
        metadata = {
            'file_name': path.name,
            'file_ext': path.suffix,
            'file_size': len(content),
            'absolute_path': str(path.absolute()),
            'paragraph_count': len(doc.paragraphs),
            'table_count': len(doc.tables),
            'document_type': 'docx'
        }

        return content, metadata

    def _parse_pptx(self, file_path: str) -> tuple:
        """
        Parse PowerPoint presentation using python-pptx

        Args:
            file_path: Path to .pptx file

        Returns:
            (content, metadata) tuple
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        path = Path(file_path)
        prs = Presentation(file_path)

        # Extract text from all slides
        slide_texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = [f"--- Slide {slide_num} ---"]

            # Extract text from all shapes in slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text)

            if len(slide_text_parts) > 1:  # More than just the header
                slide_texts.append("\n".join(slide_text_parts))

        content = "\n\n".join(slide_texts)

        # Create metadata
        metadata = {
            'file_name': path.name,
            'file_ext': path.suffix,
            'file_size': len(content),
            'absolute_path': str(path.absolute()),
            'slide_count': len(prs.slides),
            'document_type': 'pptx'
        }

        return content, metadata

    def get_stats(self) -> Dict:
        """Get loader statistics"""
        return {
            'total_loaded': self.loaded_count
        }
